from pathlib import Path

from PIL import Image, ImageDraw, ImageFilter, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(r"C:\Users\User\Documents\amazer savegarde")
OUTPUT_DIR = BASE_DIR / "output"
ASSETS_DIR = OUTPUT_DIR / "bale_assets"
OUTPUT = OUTPUT_DIR / "Expose-Comite-de-Bale-Premium-v2.pptx"
DESKTOP_OUTPUT = Path(r"C:\Users\User\Desktop\Expose-Comite-de-Bale-Premium-v2.pptx")

TITLE = "LE COMITE DE BALE"
SUBTITLE = "Regulation bancaire internationale, accords prudentiels et ratios de solvabilite"
PRESENTERS = [
    "Idrissa Tahirou Mariama",
    "Ahmed Zoubane Assalama",
    "Toudjani Ibrahim Bouchra",
    "Alzouma Nouhou Hamidou",
    "Sadek Hassane",
]


NAVY = RGBColor(15, 24, 46)
BLUE = RGBColor(31, 78, 121)
SKY = RGBColor(227, 239, 255)
ORANGE = RGBColor(237, 125, 49)
GOLD = RGBColor(255, 192, 0)
GREEN = RGBColor(42, 140, 88)
RED = RGBColor(196, 60, 60)
DARK = RGBColor(34, 34, 34)
MUTED = RGBColor(98, 103, 113)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(248, 250, 253)


def load_font(size: int, bold: bool = False):
    candidates = []
    if bold:
        candidates += [
            r"C:\Windows\Fonts\arialbd.ttf",
            r"C:\Windows\Fonts\segoeuib.ttf",
            r"C:\Windows\Fonts\calibrib.ttf",
        ]
    candidates += [
        r"C:\Windows\Fonts\arial.ttf",
        r"C:\Windows\Fonts\segoeui.ttf",
        r"C:\Windows\Fonts\calibri.ttf",
    ]
    for candidate in candidates:
        path = Path(candidate)
        if path.exists():
            return ImageFont.truetype(str(path), size=size)
    return ImageFont.load_default()


def make_cover_illustration(path: Path) -> None:
    img = Image.new("RGBA", (1400, 900), (0, 0, 0, 0))
    overlay = Image.new("RGBA", img.size, (0, 0, 0, 0))
    draw = ImageDraw.Draw(overlay)
    draw.ellipse((860, -120, 1400, 420), fill=(255, 132, 32, 210))
    draw.ellipse((720, 130, 1320, 840), fill=(50, 88, 180, 95))
    overlay = overlay.filter(ImageFilter.GaussianBlur(42))
    img.alpha_composite(overlay)
    draw = ImageDraw.Draw(img)

    # Stylized globe
    draw.ellipse((780, 90, 1220, 530), outline=(255, 255, 255, 150), width=6)
    for y in [180, 310, 440]:
        draw.arc((820, y - 80, 1180, y + 80), 180, 360, fill=(255, 255, 255, 110), width=4)
        draw.arc((820, y - 80, 1180, y + 80), 0, 180, fill=(255, 255, 255, 110), width=4)
    for x in [900, 1000, 1100]:
        draw.arc((x - 90, 110, x + 90, 510), 90, 270, fill=(255, 255, 255, 110), width=4)
        draw.arc((x - 90, 110, x + 90, 510), -90, 90, fill=(255, 255, 255, 110), width=4)

    # Bank building
    base_x, base_y = 310, 250
    draw.rounded_rectangle((base_x, base_y + 120, base_x + 360, base_y + 400), radius=18, fill=(18, 29, 58, 245))
    draw.polygon(
        [
            (base_x - 18, base_y + 120),
            (base_x + 180, base_y + 20),
            (base_x + 378, base_y + 120),
        ],
        fill=(237, 125, 49, 255),
    )
    for i in range(4):
        x = base_x + 42 + i * 72
        draw.rounded_rectangle((x, base_y + 140, x + 34, base_y + 360), radius=12, fill=(236, 240, 249, 255))
    draw.rectangle((base_x + 18, base_y + 360, base_x + 342, base_y + 392), fill=(255, 255, 255, 220))

    # Shield
    shield = [(970, 530), (1120, 530), (1160, 610), (1045, 760), (930, 610)]
    draw.polygon(shield, fill=(255, 255, 255, 235))
    draw.polygon([(992, 570), (1026, 604), (1090, 540), (1114, 564), (1026, 652), (968, 594)], fill=(42, 140, 88, 255))

    # Small cards
    for x, label in [(720, "Risque"), (900, "Capital"), (1080, "Controle")]:
        draw.rounded_rectangle((x, 640, x + 140, 710), radius=18, fill=(16, 24, 46, 220))
        draw.text((x + 26, 662), label, font=load_font(28, bold=True), fill=(255, 255, 255, 255))

    img.save(path)


def make_risk_image(path: Path) -> None:
    img = Image.new("RGBA", (1200, 700), (248, 250, 253, 255))
    draw = ImageDraw.Draw(img)
    draw.rounded_rectangle((50, 60, 1150, 640), radius=40, fill=(255, 255, 255, 255), outline=(223, 233, 247, 255), width=4)
    draw.text((80, 90), "LES QUATRE RISQUES MAJEURS", font=load_font(38, bold=True), fill=(15, 24, 46, 255))
    cards = [
        ((90, 180, 330, 520), (237, 125, 49, 255), "Credit", "Defaut de remboursement\nd'un emprunteur"),
        ((360, 180, 600, 520), (31, 78, 121, 255), "Marche", "Variation des taux,\nprix et cours"),
        ((630, 180, 870, 520), (196, 60, 60, 255), "Operationnel", "Fraude, panne,\nerreur humaine"),
        ((900, 180, 1140, 520), (42, 140, 88, 255), "Liquidite", "Capacite a honorer\nles engagements"),
    ]
    for box, color, title, body in cards:
        draw.rounded_rectangle(box, radius=28, fill=(255, 255, 255, 255), outline=color, width=6)
        x1, y1, x2, y2 = box
        draw.ellipse((x1 + 72, y1 + 26, x1 + 148, y1 + 102), fill=color)
        draw.text((x1 + 30, y1 + 126), title, font=load_font(30, bold=True), fill=(15, 24, 46, 255))
        draw.multiline_text((x1 + 30, y1 + 190), body, font=load_font(22), fill=(70, 74, 82, 255), spacing=8)
    img.save(path)


def make_ratios_dashboard(path: Path) -> None:
    img = Image.new("RGBA", (1200, 760), (255, 255, 255, 255))
    draw = ImageDraw.Draw(img)
    draw.rounded_rectangle((40, 40, 1160, 720), radius=36, fill=(15, 24, 46, 255))
    draw.text((90, 86), "DASHBOARD DE SOLVABILITE", font=load_font(38, bold=True), fill=(255, 255, 255, 255))

    # cards
    cards = [
        ((90, 170, 400, 340), "Ratio Cooke", ">= 8%", (237, 125, 49, 255)),
        ((430, 170, 740, 340), "Ratio de levier", ">= 3%", (255, 192, 0, 255)),
        ((770, 170, 1080, 340), "Capital CET1", "Qualite renforcee", (42, 140, 88, 255)),
    ]
    for box, title, value, color in cards:
        draw.rounded_rectangle(box, radius=26, fill=(255, 255, 255, 245))
        x1, y1, x2, y2 = box
        draw.rounded_rectangle((x1 + 22, y1 + 18, x1 + 126, y1 + 54), radius=16, fill=color)
        draw.text((x1 + 34, y1 + 24), title, font=load_font(18, bold=True), fill=(255, 255, 255, 255))
        draw.text((x1 + 24, y1 + 92), value, font=load_font(34, bold=True), fill=(15, 24, 46, 255))

    # mini bars
    draw.rounded_rectangle((90, 400, 1080, 650), radius=28, fill=(255, 255, 255, 245))
    draw.text((120, 430), "Lecture simple des seuils prudentiels", font=load_font(28, bold=True), fill=(15, 24, 46, 255))
    labels = [("Cooke", 0.78, (237, 125, 49, 255)), ("Levier", 0.38, (255, 192, 0, 255)), ("Liquidite", 0.66, (42, 140, 88, 255))]
    y = 495
    for label, ratio, color in labels:
        draw.text((120, y - 10), label, font=load_font(22, bold=True), fill=(60, 65, 74, 255))
        draw.rounded_rectangle((260, y, 910, y + 28), radius=14, fill=(228, 234, 243, 255))
        draw.rounded_rectangle((260, y, int(260 + 650 * ratio), y + 28), radius=14, fill=color)
        draw.text((940, y - 8), f"{int(ratio * 100)}%", font=load_font(22, bold=True), fill=(15, 24, 46, 255))
        y += 62
    img.save(path)


def make_pillars_image(path: Path) -> None:
    img = Image.new("RGBA", (1200, 760), (248, 250, 253, 255))
    draw = ImageDraw.Draw(img)
    draw.rounded_rectangle((42, 42, 1158, 718), radius=36, fill=(255, 255, 255, 255), outline=(223, 233, 247, 255), width=4)
    draw.text((88, 82), "LES 3 PILIERS DE BALE II", font=load_font(38, bold=True), fill=(15, 24, 46, 255))
    colors = [(31, 78, 121, 255), (237, 125, 49, 255), (42, 140, 88, 255)]
    titles = ["Pilier 1", "Pilier 2", "Pilier 3"]
    subtitles = [
        "Exigences minimales\nen fonds propres",
        "Surveillance\nprudentielle",
        "Discipline\nde marche",
    ]
    descs = [
        "Mesurer les risques et exiger\nun capital adequat.",
        "Verifier que les banques\ngerent correctement leurs risques.",
        "Renforcer la transparence\ndes informations publiees.",
    ]
    xs = [90, 400, 710]
    for x, color, title, subtitle, desc in zip(xs, colors, titles, subtitles, descs):
        draw.rounded_rectangle((x, 180, x + 250, 560), radius=30, fill=(255, 255, 255, 255), outline=color, width=6)
        draw.rounded_rectangle((x + 22, 202, x + 228, 255), radius=18, fill=color)
        draw.text((x + 70, 215), title, font=load_font(24, bold=True), fill=(255, 255, 255, 255))
        draw.multiline_text((x + 25, 292), subtitle, font=load_font(28, bold=True), fill=(15, 24, 46, 255), spacing=8)
        draw.multiline_text((x + 25, 405), desc, font=load_font(20), fill=(85, 90, 100, 255), spacing=8)
    draw.line((340, 370, 400, 370), fill=(170, 180, 198, 255), width=6)
    draw.line((650, 370, 710, 370), fill=(170, 180, 198, 255), width=6)
    img.save(path)


def make_stability_image(path: Path) -> None:
    img = Image.new("RGBA", (1200, 760), (255, 255, 255, 255))
    overlay = Image.new("RGBA", img.size, (0, 0, 0, 0))
    odraw = ImageDraw.Draw(overlay)
    odraw.ellipse((760, 40, 1190, 470), fill=(255, 140, 30, 120))
    odraw.ellipse((690, 320, 1180, 760), fill=(31, 78, 121, 55))
    overlay = overlay.filter(ImageFilter.GaussianBlur(40))
    img.alpha_composite(overlay)
    draw = ImageDraw.Draw(img)
    draw.text((80, 88), "DE LA REGLE A LA STABILITE", font=load_font(38, bold=True), fill=(15, 24, 46, 255))

    stages = [
        ("Risque bancaire", (196, 60, 60, 255)),
        ("Regles de Bale", (237, 125, 49, 255)),
        ("Banque plus solide", (31, 78, 121, 255)),
        ("Stabilite financiere", (42, 140, 88, 255)),
    ]
    xs = [90, 350, 610, 870]
    for i, ((label, color), x) in enumerate(zip(stages, xs)):
        draw.rounded_rectangle((x, 300, x + 220, 430), radius=28, fill=(255, 255, 255, 245), outline=color, width=5)
        draw.ellipse((x + 74, 220, x + 146, 292), fill=color)
        bbox = draw.textbbox((0, 0), label, font=load_font(24, bold=True))
        tw = bbox[2] - bbox[0]
        draw.text((x + (220 - tw) / 2, 346), label, font=load_font(24, bold=True), fill=(15, 24, 46, 255))
        if i < len(xs) - 1:
            draw.line((x + 220, 365, xs[i + 1], 365), fill=(167, 176, 192, 255), width=7)
            draw.polygon([(xs[i + 1] - 18, 355), (xs[i + 1], 365), (xs[i + 1] - 18, 375)], fill=(167, 176, 192, 255))
    img.save(path)


def make_assets():
    ASSETS_DIR.mkdir(parents=True, exist_ok=True)
    cover = ASSETS_DIR / "cover.png"
    risks = ASSETS_DIR / "risks.png"
    ratios = ASSETS_DIR / "ratios.png"
    pillars = ASSETS_DIR / "pillars.png"
    stability = ASSETS_DIR / "stability.png"
    make_cover_illustration(cover)
    make_risk_image(risks)
    make_ratios_dashboard(ratios)
    make_pillars_image(pillars)
    make_stability_image(stability)
    return cover, risks, ratios, pillars, stability


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_text(slide, x, y, w, h, text, size=28, color=NAVY, bold=True):
    box = slide.shapes.add_textbox(x, y, w, h)
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_body_bullets(slide, x, y, w, h, bullets, size=19, color=DARK):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for idx, item in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.level = 0
        p.bullet = True
    return box


def add_top_band(slide, title, section):
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.72))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    add_title_text(slide, Inches(0.5), Inches(0.15), Inches(7), Inches(0.35), title, size=24, color=WHITE)
    pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(10.45), Inches(0.12), Inches(2.3), Inches(0.4))
    pill.fill.solid()
    pill.fill.fore_color.rgb = ORANGE
    pill.line.fill.background()
    add_title_text(slide, Inches(10.63), Inches(0.18), Inches(1.95), Inches(0.25), section, size=13, color=WHITE)


def add_footer(slide, right="Comité de Bâle"):
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(7.03), Inches(13.333), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = SKY
    line.line.fill.background()
    add_title_text(slide, Inches(0.45), Inches(7.08), Inches(4.8), Inches(0.2), "Université Islamique au Niger", size=10, color=MUTED, bold=False)
    box = slide.shapes.add_textbox(Inches(10.1), Inches(7.08), Inches(2.6), Inches(0.2))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = right
    r.font.size = Pt(10)
    r.font.color.rgb = MUTED


def build_presentation():
    cover_img, risks_img, ratios_img, pillars_img, stability_img = make_assets()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    top_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.18))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = ORANGE
    top_bar.line.fill.background()
    slide.shapes.add_picture(str(cover_img), Inches(7.1), Inches(0.55), width=Inches(5.6))
    add_title_text(slide, Inches(0.62), Inches(0.82), Inches(5.6), Inches(0.6), TITLE, size=28, color=ORANGE)
    add_title_text(slide, Inches(0.62), Inches(1.5), Inches(5.8), Inches(1.8), "Une lecture claire\net visuelle de la\nrégulation bancaire", size=30, color=NAVY)
    add_title_text(slide, Inches(0.65), Inches(3.9), Inches(5.8), Inches(0.6), SUBTITLE, size=18, color=MUTED, bold=False)

    names_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(4.65), Inches(5.7), Inches(1.8))
    names_box.fill.solid()
    names_box.fill.fore_color.rgb = WHITE
    names_box.line.color.rgb = SKY
    tf = names_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Présenté par"
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = BLUE
    for name in PRESENTERS:
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = f"• {name}"
        r.font.size = Pt(15)
        r.font.color.rgb = DARK
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = "Professeur : Mr Ahmad Sidi Ibrahim Mazou | 2025-2026"
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = MUTED
    add_footer(slide, right="Exposé académique")

    # 2. Plan cards
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_top_band(slide, "Plan de l'exposé", "Plan")
    plan_items = [
        ("01", "Introduction", "Définition et utilité du Comité de Bâle"),
        ("02", "Présentation", "Création, objectifs et rôle"),
        ("03", "Principes", "Fonds propres, risques, contrôle interne"),
        ("04", "Accords", "Bâle I, II et III"),
        ("05", "Piliers", "Les 3 piliers de Bâle II"),
        ("06", "Ratios", "Cooke et levier"),
    ]
    x_positions = [0.7, 4.45, 8.2]
    y_positions = [1.3, 4.0]
    idx = 0
    for y in y_positions:
        for x in x_positions:
            num, title, desc = plan_items[idx]
            card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.1), Inches(2.2))
            card.fill.solid()
            card.fill.fore_color.rgb = LIGHT
            card.line.color.rgb = SKY
            add_title_text(slide, Inches(x + 0.22), Inches(y + 0.18), Inches(0.6), Inches(0.3), num, size=24, color=ORANGE)
            add_title_text(slide, Inches(x + 0.22), Inches(y + 0.62), Inches(2.45), Inches(0.35), title, size=21, color=NAVY)
            add_title_text(slide, Inches(x + 0.22), Inches(y + 1.1), Inches(2.55), Inches(0.72), desc, size=15, color=MUTED, bold=False)
            idx += 1
    add_footer(slide)

    # 3. Intro with image
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_top_band(slide, "Introduction et définition", "Introduction")
    slide.shapes.add_picture(str(cover_img), Inches(7.45), Inches(1.15), width=Inches(4.95))
    add_title_text(slide, Inches(0.7), Inches(1.25), Inches(5.8), Inches(0.4), "Pourquoi le Comité de Bâle existe-t-il ?", size=25, color=NAVY)
    add_body_bullets(
        slide,
        Inches(0.82),
        Inches(1.9),
        Inches(5.7),
        Inches(3.9),
        [
            "Les banques sont interconnectées à l'échelle mondiale.",
            "La faillite d'une grande banque peut provoquer un effet de contagion sur plusieurs pays.",
            "Le Comité de Bâle fixe des standards prudentiels pour limiter ces crises.",
            "Il regroupe des banques centrales et autorités de surveillance, sans voter lui-même des lois.",
        ],
    )
    quote = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.82), Inches(5.85), Inches(5.85), Inches(0.9))
    quote.fill.solid()
    quote.fill.fore_color.rgb = NAVY
    quote.line.fill.background()
    add_title_text(slide, Inches(1.05), Inches(6.03), Inches(5.2), Inches(0.45), "En bref : coordonner la sécurité bancaire internationale.", size=19, color=WHITE)
    add_footer(slide)

    # 4. Creation, objectifs, role
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_top_band(slide, "Création, objectifs et rôle", "Présentation")
    add_title_text(slide, Inches(0.72), Inches(1.0), Inches(3), Inches(0.3), "Création en 1974", size=23, color=NAVY)
    add_title_text(slide, Inches(0.72), Inches(1.45), Inches(3.3), Inches(0.7), "Suite à la faillite de la banque Herstatt, qui a montré la nécessité de règles communes à l'international.", size=17, color=MUTED, bold=False)
    cards = [
        (Inches(0.72), "Objectif", "Rendre les banques plus solides et plus résilientes."),
        (Inches(4.35), "Rôle", "Fixer des standards, surveiller les risques et harmoniser la coopération."),
        (Inches(8.0), "Finalité", "Prévenir les crises et éviter des sauvetages coûteux pour les États."),
    ]
    for x, title, desc in cards:
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, Inches(2.35), Inches(3.1), Inches(2.8))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT
        card.line.color.rgb = SKY
        badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x + Inches(0.22), Inches(2.58), Inches(0.55), Inches(0.55))
        badge.fill.solid()
        badge.fill.fore_color.rgb = ORANGE
        badge.line.fill.background()
        add_title_text(slide, x + Inches(0.92), Inches(2.63), Inches(1.8), Inches(0.25), title, size=20, color=NAVY)
        add_title_text(slide, x + Inches(0.28), Inches(3.28), Inches(2.55), Inches(1.2), desc, size=16, color=MUTED, bold=False)
    add_footer(slide)

    # 5. Principles + image
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_top_band(slide, "Les principes de base", "Principes")
    slide.shapes.add_picture(str(risks_img), Inches(7.15), Inches(1.2), width=Inches(5.2))
    principle_boxes = [
        ("Fonds propres", "Capital minimum pour absorber les pertes imprévues."),
        ("Gestion des risques", "Crédit, marché, opérationnel et liquidité doivent être identifiés et mesurés."),
        ("Contrôle interne", "Surveillance, conformité, gouvernance et fiabilité de l'information financière."),
    ]
    y = 1.35
    for title, desc in principle_boxes:
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.78), Inches(y), Inches(5.7), Inches(1.35))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = SKY
        add_title_text(slide, Inches(1.05), Inches(y + 0.18), Inches(2.8), Inches(0.28), title, size=20, color=BLUE)
        add_title_text(slide, Inches(1.05), Inches(y + 0.56), Inches(4.85), Inches(0.46), desc, size=15, color=MUTED, bold=False)
        y += 1.63
    add_footer(slide)

    # 6. Accords comparison table
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_top_band(slide, "Comparaison des accords de Bâle", "Tableau")
    table = slide.shapes.add_table(4, 4, Inches(0.6), Inches(1.3), Inches(12.0), Inches(4.0)).table
    headers = ["Accord", "Année", "Apport majeur", "Limite / évolution"]
    rows = [
        ["Bâle I", "1988", "Ratio Cooke, seuil minimal de 8% des actifs pondérés par le risque", "Vision trop simpliste des risques"],
        ["Bâle II", "2004", "Trois piliers, meilleure gestion prudentielle et discipline de marché", "Insuffisant face à la crise de 2008"],
        ["Bâle III", "2010", "Capital renforcé, liquidité, ratio de levier", "Réponse corrective aux faiblesses précédentes"],
    ]
    for i, head in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = head
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(17)
        para.font.bold = True
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if r_idx % 2 == 1 else SKY
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(15)
            para.font.color.rgb = DARK
            if c_idx in (0, 1):
                para.font.bold = True
                para.alignment = PP_ALIGN.CENTER
    info = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(5.75), Inches(11.5), Inches(0.72))
    info.fill.solid()
    info.fill.fore_color.rgb = ORANGE
    info.line.fill.background()
    add_title_text(slide, Inches(1.2), Inches(5.96), Inches(11), Inches(0.28), "Lecture simple : Bâle I pose les bases, Bâle II approfondit, Bâle III renforce durablement.", size=18, color=WHITE)
    add_footer(slide)

    # 7. Bale II pillars
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_top_band(slide, "Les trois piliers de Bâle II", "Bâle II")
    slide.shapes.add_picture(str(pillars_img), Inches(6.7), Inches(1.1), width=Inches(5.55))
    add_title_text(slide, Inches(0.72), Inches(1.2), Inches(5.4), Inches(0.35), "Une logique prudentielle en 3 dimensions", size=24, color=NAVY)
    add_body_bullets(
        slide,
        Inches(0.85),
        Inches(1.8),
        Inches(5.35),
        Inches(3.8),
        [
            "Pilier 1 : exigences minimales en fonds propres.",
            "Pilier 2 : surveillance prudentielle par les autorités.",
            "Pilier 3 : discipline de marché grâce à la transparence.",
            "Bâle II cherche à mieux refléter la réalité des risques bancaires.",
        ],
        size=18,
    )
    key = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.82), Inches(5.78), Inches(5.5), Inches(0.82))
    key.fill.solid()
    key.fill.fore_color.rgb = WHITE
    key.line.color.rgb = SKY
    add_title_text(slide, Inches(1.05), Inches(6.01), Inches(5.0), Inches(0.25), "Idée clé : plus de risque = plus d'exigence prudentielle.", size=17, color=BLUE)
    add_footer(slide)

    # 8. Timeline
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_top_band(slide, "Frise d'évolution des accords", "Chronologie")
    base_line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.05), Inches(3.35), Inches(11.0), Inches(0.08))
    base_line.fill.solid()
    base_line.fill.fore_color.rgb = BLUE
    base_line.line.fill.background()
    timeline = [
        (1.4, "1974", "Création", "Naissance du Comité après la faillite Herstatt"),
        (4.2, "1988", "Bâle I", "Premier cadre de solvabilité"),
        (7.0, "2004", "Bâle II", "Trois piliers prudentiels"),
        (9.8, "2010", "Bâle III", "Capital, liquidité et levier"),
    ]
    for x, year, title, desc in timeline:
        circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(3.05), Inches(0.45), Inches(0.45))
        circ.fill.solid()
        circ.fill.fore_color.rgb = ORANGE
        circ.line.fill.background()
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x - 0.55), Inches(1.65), Inches(2.2), Inches(1.15))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = SKY
        add_title_text(slide, Inches(x - 0.32), Inches(1.8), Inches(1.4), Inches(0.22), year, size=18, color=ORANGE)
        add_title_text(slide, Inches(x - 0.32), Inches(2.05), Inches(1.6), Inches(0.22), title, size=18, color=NAVY)
        add_title_text(slide, Inches(x - 0.32), Inches(2.36), Inches(1.75), Inches(0.38), desc, size=11, color=MUTED, bold=False)
    add_footer(slide)

    # 9. Bale III innovations
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_top_band(slide, "Bâle III : les innovations majeures", "Bâle III")
    cards = [
        ("Capital de meilleure qualité", "Accent sur les fonds propres les plus solides, notamment le CET1.", ORANGE),
        ("Exigences de liquidité", "Les banques doivent pouvoir faire face à leurs engagements à court terme.", BLUE),
        ("Ratio de levier", "Levier minimal de 3% pour limiter l'endettement excessif.", GREEN),
    ]
    xs = [0.7, 4.45, 8.2]
    for (title, desc, color), x in zip(cards, xs):
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.65), Inches(3.1), Inches(3.2))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT
        card.line.color.rgb = color
        chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + 0.22), Inches(1.88), Inches(1.2), Inches(0.32))
        chip.fill.solid()
        chip.fill.fore_color.rgb = color
        chip.line.fill.background()
        add_title_text(slide, Inches(x + 0.4), Inches(1.95), Inches(0.85), Inches(0.18), "Bâle III", size=11, color=WHITE)
        add_title_text(slide, Inches(x + 0.22), Inches(2.45), Inches(2.55), Inches(0.65), title, size=20, color=NAVY)
        add_title_text(slide, Inches(x + 0.22), Inches(3.35), Inches(2.45), Inches(0.95), desc, size=15, color=MUTED, bold=False)
    slide.shapes.add_picture(str(stability_img), Inches(1.25), Inches(5.15), width=Inches(10.9))
    add_footer(slide)

    # 10. Ratios with dashboard image
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_top_band(slide, "Les principaux ratios de solvabilité", "Ratios")
    slide.shapes.add_picture(str(ratios_img), Inches(7.12), Inches(1.25), width=Inches(5.1))
    for y, title, formula, threshold, note in [
        (1.3, "Ratio de Cooke", "Fonds propres / Actifs pondérés par les risques x 100", "Seuil : 8%", "Mesure la capacité à couvrir les pertes sur les actifs risqués."),
        (3.15, "Ratio de levier", "Fonds propres de base / Total des expositions x 100", "Seuil : 3%", "Empêche un endettement excessif, même lorsque le risque paraît faible."),
    ]:
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(y), Inches(5.85), Inches(1.5))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT
        card.line.color.rgb = SKY
        add_title_text(slide, Inches(0.98), Inches(y + 0.14), Inches(2.3), Inches(0.24), title, size=20, color=BLUE)
        add_title_text(slide, Inches(0.98), Inches(y + 0.48), Inches(5.1), Inches(0.25), formula, size=14, color=DARK, bold=False)
        add_title_text(slide, Inches(0.98), Inches(y + 0.82), Inches(1.4), Inches(0.25), threshold, size=16, color=ORANGE)
        add_title_text(slide, Inches(2.2), Inches(y + 0.84), Inches(4.0), Inches(0.25), note, size=13, color=MUTED, bold=False)
    add_footer(slide)

    # 11. Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    add_title_text(slide, Inches(0.7), Inches(0.75), Inches(4.8), Inches(0.4), "Conclusion", size=28, color=WHITE)
    summary = [
        "Le Comité de Bâle encadre la sécurité du système bancaire international.",
        "Les accords de Bâle I, II et III ont progressivement renforcé la solvabilité des banques.",
        "Les ratios prudentiels permettent d'évaluer la solidité financière et de limiter les crises.",
        "L'enjeu final est la stabilité financière et la protection de l'économie réelle.",
    ]
    y = 1.55
    for item in summary:
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.78), Inches(y), Inches(8.0), Inches(0.85))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(23, 38, 68)
        card.line.color.rgb = RGBColor(58, 84, 130)
        add_title_text(slide, Inches(1.02), Inches(y + 0.2), Inches(7.3), Inches(0.25), f"• {item}", size=17, color=WHITE)
        y += 1.0
    thanks = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.2), Inches(1.7), Inches(3.15), Inches(2.2))
    thanks.fill.solid()
    thanks.fill.fore_color.rgb = ORANGE
    thanks.line.fill.background()
    box = slide.shapes.add_textbox(Inches(9.45), Inches(2.15), Inches(2.6), Inches(1.1))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "MERCI\nPOUR VOTRE ATTENTION"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = WHITE
    add_footer(slide, right="Questions / Réponses")

    return prs


def main():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = build_presentation()
    prs.save(OUTPUT)
    prs.save(DESKTOP_OUTPUT)
    print(OUTPUT)
    print(DESKTOP_OUTPUT)


if __name__ == "__main__":
    main()
